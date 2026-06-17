from __future__ import annotations

from pathlib import Path

import pytest

from src.slide_export import auto_deck
from src import cli
from src.slide_export.slide_manifest import generate_slide_manifest
from src.slide_export.pptx_generator import DECK_FONT_FAMILY, DELIVEROO_COLOR, GLOVO_COLOR, SECTION_TITLE_SIZE_SHORT, generate_slides
from src.slide_export.template_variants import resolve_template_id


def test_generate_slide_manifest_creates_manifest() -> None:
    generate_slide_manifest()

    manifest = Path("outputs/slide_manifest.md")
    assert manifest.exists()
    assert "Slide Manifest" in manifest.read_text(encoding="utf-8")


def test_generate_slides_from_yaml_template(tmp_path: Path) -> None:
    from PIL import Image
    from pptx import Presentation

    template = tmp_path / "template.pptx"
    config = tmp_path / "deck.yml"
    image = tmp_path / "graph.png"
    table = tmp_path / "table.csv"
    output = tmp_path / "out.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(100000, 100000, 3000000, 300000)
    box.text = "TEMPLATE_ID:cover"
    slide.shapes.add_textbox(100000, 500000, 5000000, 500000).text = "[PROJECT_TITLE]"
    slide.shapes.add_textbox(100000, 1100000, 3000000, 1800000).text = "[GRAPH_MAIN]"
    slide.shapes.add_textbox(100000, 3000000, 3000000, 1000000).text = "[TABLE_MAIN]"
    prs.save(template)

    Image.new("RGB", (640, 360), "#00CCBC").save(image)
    table.write_text("Metric,Value\nSuccess,92%\n", encoding="utf-8")
    config.write_text(
        f"""
deck:
  template: "{template.as_posix()}"
  output: "{output.as_posix()}"
slides:
  - template_id: cover
    fields:
      PROJECT_TITLE: "Deck demo"
    images:
      GRAPH_MAIN: "{image.as_posix()}"
    table:
      placeholder: TABLE_MAIN
      source: "{table.as_posix()}"
""",
        encoding="utf-8",
    )

    result = generate_slides(config, overwrite=True)

    assert result.generated_slides == 1
    assert output.exists()
    generated = Presentation(output)
    assert len(generated.slides) == 1
    assert "Deck demo" in "\n".join(shape.text for shape in generated.slides[0].shapes if getattr(shape, "has_text_frame", False))


def test_generate_slides_replaces_text_by_shape_name(tmp_path: Path) -> None:
    from pptx import Presentation

    template = tmp_path / "template.pptx"
    config = tmp_path / "deck.yml"
    output = tmp_path / "out.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(100000, 100000, 3000000, 300000).text = "TEMPLATE_ID:sources"
    source_box = slide.shapes.add_textbox(100000, 500000, 5000000, 500000)
    source_box.name = "SOURCES_LIST"
    source_box.text = "Old template source"
    prs.save(template)

    config.write_text(
        f"""
deck:
  template: "{template.as_posix()}"
  output: "{output.as_posix()}"
slides:
  - template_id: sources
    fields:
      SOURCES_LIST: "outputs/figures/dark/example.png"
""",
        encoding="utf-8",
    )

    generate_slides(config, overwrite=True)

    generated = Presentation(output)
    text = "\n".join(shape.text for shape in generated.slides[0].shapes if getattr(shape, "has_text_frame", False))
    assert "outputs/figures/dark/example.png" in text
    assert "Old template source" not in text


def test_auto_slide_expansion_uses_existing_assets(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    def resolve_in_tmp(path: str | Path, base_dir: Path = tmp_path) -> Path:
        value = Path(path)
        return value if value.is_absolute() else tmp_path / value

    monkeypatch.setattr(auto_deck, "resolve_path", resolve_in_tmp)
    figure = tmp_path / "outputs/figures/dark/questionnaire/ueq_scales.png"
    table = tmp_path / "outputs/tables/ueq_summary.csv"
    snippet = tmp_path / "outputs/texts/snippets/questionnaire_conclusions.md"
    figure.parent.mkdir(parents=True)
    table.parent.mkdir(parents=True)
    snippet.parent.mkdir(parents=True)
    figure.write_bytes(b"fake-png")
    table.write_text("metric,value\n" + "\n".join(f"score_{idx},{idx}" for idx in range(15)) + "\n", encoding="utf-8")
    snippet.write_text("# Conclusioni\n\n- La differenza principale emerge negli item di chiarezza e controllo.\n", encoding="utf-8")

    expanded = auto_deck.expand_auto_slides(
        {"auto_slides": {"enabled": True}, "slides": []},
        {"graph_full", "table_large", "findings", "section_divider", "sources"},
    )

    template_ids = [slide.get("template_id") or slide.get("template") for slide in expanded["slides"]]
    assert "graph_full" in template_ids
    assert template_ids.count("table_large") == 2
    assert "findings" in template_ids
    assert "sources" in template_ids


def test_auto_slide_expansion_creates_app_variants(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    def resolve_in_tmp(path: str | Path, base_dir: Path = tmp_path) -> Path:
        value = Path(path)
        return value if value.is_absolute() else tmp_path / value

    monkeypatch.setattr(auto_deck, "resolve_path", resolve_in_tmp)
    task_figure = tmp_path / "outputs/figures/dark/user_tests/tasks/t01_effectiveness.png"
    ueq_figure = tmp_path / "outputs/figures/dark/questionnaire/items/item_01_boxplot.png"
    task_table = tmp_path / "outputs/tables/user_tests_t01_summary.csv"
    ueq_table = tmp_path / "outputs/tables/questionnaire_items_summary.csv"
    text = tmp_path / "outputs/texts/snippets/nps.md"
    for path in [task_figure, ueq_figure, task_table, ueq_table, text]:
        path.parent.mkdir(parents=True, exist_ok=True)
    task_figure.write_bytes(b"fake-png")
    ueq_figure.write_bytes(b"fake-png")
    task_table.write_text("app,success_rate,mean_time_sec\nDeliveroo,0.8,22\nGlovo,1.0,18\n", encoding="utf-8")
    ueq_table.write_text(
        "item_number,item,Deliveroo_mean,Deliveroo_median,Glovo_mean,Glovo_median,mean_difference_abs,p_value,median_difference_abs\n"
        "1,Chiarezza,1.2,1,0.6,0,0.6,0.04,1\n",
        encoding="utf-8",
    )
    text.write_text("# NPS\n\nDeliveroo risulta preferibile nel consiglio agli utenti, con un vantaggio chiaro nel dato sintetico.", encoding="utf-8")

    expanded = auto_deck.expand_auto_slides(
        {"auto_slides": {"enabled": True}, "slides": []},
        {
            "section_divider_deliveroo",
            "section_divider_glovo",
            "task_results_deliveroo",
            "task_results_glovo",
            "ueq_question_deliveroo",
            "ueq_question_glovo",
            "findings_deliveroo",
            "sources",
        },
    )

    themed = [(slide.get("template"), slide.get("theme")) for slide in expanded["slides"]]
    assert ("task_results", "deliveroo") in themed
    assert ("task_results", "glovo") in themed
    assert ("ueq_question", "deliveroo") in themed
    assert ("ueq_question", "glovo") in themed
    assert ("findings", "deliveroo") in themed


def test_reference_order_slide_expansion_follows_reference_sections(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    def resolve_in_tmp(path: str | Path, base_dir: Path = tmp_path) -> Path:
        value = Path(path)
        return value if value.is_absolute() else tmp_path / value

    monkeypatch.setattr(auto_deck, "resolve_path", resolve_in_tmp)
    content = tmp_path / "slides/content/reference_static_texts.md"
    content.parent.mkdir(parents=True)
    content.write_text(
        "# Testi\n\n"
        "## problem_description\nProblema delivery.\n\n"
        "## app_deliveroo\nTesto Deliveroo.\n\n"
        "## app_glovo\nTesto Glovo.\n",
        encoding="utf-8",
    )
    figure = tmp_path / "outputs/figures/dark/questionnaire/ueq_scales.png"
    expertise = tmp_path / "outputs/figures/dark/heuristics/expertise_matrix.png"
    errors = tmp_path / "outputs/figures/dark/user_tests/tasks/t01_error_breakdown.png"
    success = tmp_path / "outputs/figures/dark/users_time_success_rate.png"
    boxplot = tmp_path / "outputs/figures/dark/users_time_boxplot_by_task.png"
    figure.parent.mkdir(parents=True)
    expertise.parent.mkdir(parents=True, exist_ok=True)
    errors.parent.mkdir(parents=True, exist_ok=True)
    for path in [figure, expertise, errors, success, boxplot]:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(b"fake-png")

    expanded = auto_deck.expand_auto_slides(
        {
            "metadata": {"project_title": "Deliveroo vs Glovo", "subtitle": "Analisi trasversale", "authors": "Team", "date": "2026"},
            "auto_slides": {"enabled": True, "mode": "reference_order", "reference_texts": "slides/content/reference_static_texts.md"},
            "slides": [{"template_id": "cover", "fields": {"PROJECT_TITLE": "Manuale"}}],
        },
        {
            "cover",
            "section_divider_neutral",
            "text_only_neutral",
            "text_only_deliveroo",
            "text_only_glovo",
            "graph_full_neutral",
            "comparison_neutral",
        },
    )

    slides = expanded["slides"]
    titles = [slide.get("fields", {}).get("TEXT_TITLE") or slide.get("fields", {}).get("SECTION_NAME") for slide in slides[:10]]
    assert slides[0]["template_id"] == "cover"
    assert titles[1:8] == [
        "Indice",
        "Introduzione",
        "Descrizione del problema",
        "Ambiente di valutazione",
        "Deliveroo",
        "Glovo",
        "Valutazione euristica",
    ]
    assert any(slide.get("fields", {}).get("GRAPH_TITLE") == "Media risultati UEQ" for slide in slides)
    assert any(slide.get("fields", {}).get("GRAPH_TITLE") == "Matrice di expertise" for slide in slides)
    assert any(slide.get("fields", {}).get("GRAPH_TITLE") == "Errori - Task 1" for slide in slides)
    assert any(slide.get("fields", {}).get("COMPARISON_TITLE") == "Successo e distribuzione tempi" for slide in slides)
    assert all(slide.get("fields", {}).get("PROJECT_TITLE") != "Manuale" for slide in slides)


def test_user_task_deck_expansion_creates_participant_deck(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    def resolve_in_tmp(path: str | Path, base_dir: Path = tmp_path) -> Path:
        value = Path(path)
        return value if value.is_absolute() else tmp_path / value

    monkeypatch.setattr(auto_deck, "resolve_path", resolve_in_tmp)
    content = tmp_path / "slides/content/reference_static_texts.md"
    content.parent.mkdir(parents=True)
    content.write_text(
        "# Testi\n\n"
        "## task_deck_purpose\nScopo del test.\n\n"
        "## task_deck_before_start\nPrima di iniziare.\n\n"
        "## task_deck_survey\nLink survey.\n\n"
        + "\n\n".join(
            f"## user_task_{idx}_{theme}\nTask {idx} {theme}."
            for theme in ["deliveroo", "glovo"]
            for idx in range(1, 6)
        ),
        encoding="utf-8",
    )

    expanded = auto_deck.expand_auto_slides(
        {
            "metadata": {"project_title": "User test Deliveroo vs Glovo", "subtitle": "Task", "authors": "Team", "date": "2026"},
            "auto_slides": {"enabled": True, "mode": "user_tasks", "reference_texts": "slides/content/reference_static_texts.md"},
            "slides": [],
        },
        {"cover", "section_divider_neutral", "text_only_neutral", "text_only_deliveroo", "text_only_glovo"},
    )

    slides = expanded["slides"]
    titles = [slide.get("fields", {}).get("TEXT_TITLE") or slide.get("fields", {}).get("SECTION_NAME") or slide.get("fields", {}).get("PROJECT_TITLE") for slide in slides]
    assert len(slides) == 16
    assert titles[:4] == ["User test Deliveroo vs Glovo", "A cosa serve", "Prima di iniziare", "Deliveroo"]
    assert slides[3]["theme"] == "deliveroo"
    assert titles[4:9] == [f"Deliveroo - Task {idx}" for idx in range(1, 6)]
    assert titles[9] == "Glovo"
    assert slides[9]["theme"] == "glovo"
    assert titles[10:15] == [f"Glovo - Task {idx}" for idx in range(1, 6)]
    assert titles[15] == "Conclusione e questionario"


def test_user_task_deck_respects_configured_task_count(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    def resolve_in_tmp(path: str | Path, base_dir: Path = tmp_path) -> Path:
        value = Path(path)
        return value if value.is_absolute() else tmp_path / value

    monkeypatch.setattr(auto_deck, "resolve_path", resolve_in_tmp)
    content = tmp_path / "slides/content/reference_static_texts.md"
    content.parent.mkdir(parents=True)
    content.write_text(
        "# Testi\n\n"
        "## task_deck_purpose\nScopo del test.\n\n"
        "## task_deck_before_start\nPrima di iniziare.\n\n"
        "## task_deck_survey\nLink survey.\n\n"
        + "\n\n".join(
            f"## user_task_{idx}_{theme}\nTask {idx} {theme}."
            for theme in ["deliveroo", "glovo"]
            for idx in range(1, 4)
        ),
        encoding="utf-8",
    )

    expanded = auto_deck.expand_auto_slides(
        {
            "metadata": {"project_title": "User test Deliveroo vs Glovo", "subtitle": "Task", "authors": "Team", "date": "2026"},
            "auto_slides": {
                "enabled": True,
                "mode": "user_tasks",
                "reference_texts": "slides/content/reference_static_texts.md",
                "task_count": 3,
            },
            "slides": [],
        },
        {"cover", "section_divider_neutral", "text_only_neutral", "text_only_deliveroo", "text_only_glovo"},
    )

    titles = [slide.get("fields", {}).get("TEXT_TITLE") or slide.get("fields", {}).get("SECTION_NAME") or slide.get("fields", {}).get("PROJECT_TITLE") for slide in expanded["slides"]]
    assert len(titles) == 12
    assert titles[4:7] == [f"Deliveroo - Task {idx}" for idx in range(1, 4)]
    assert titles[8:11] == [f"Glovo - Task {idx}" for idx in range(1, 4)]
    assert expanded["slides"][3]["theme"] == "deliveroo"
    assert expanded["slides"][7]["theme"] == "glovo"


def test_full_pipeline_generates_final_and_user_task_decks(monkeypatch: pytest.MonkeyPatch) -> None:
    calls = []

    def fake_generate_slides(config_path: str, *, overwrite: bool = False, timestamp: bool = False):
        calls.append((config_path, overwrite, timestamp))
        return config_path

    monkeypatch.setattr(cli, "generate_slides", fake_generate_slides)

    results = cli.generate_full_pipeline_slide_decks(overwrite=True, timestamp=True)

    assert results == ["slides/config/slide_deck.yml"]
    assert calls == [("slides/config/slide_deck.yml", True, True)]


def test_reference_order_adds_appendix_placeholders_without_assets(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    def resolve_in_tmp(path: str | Path, base_dir: Path = tmp_path) -> Path:
        value = Path(path)
        return value if value.is_absolute() else tmp_path / value

    monkeypatch.setattr(auto_deck, "resolve_path", resolve_in_tmp)
    content = tmp_path / "slides/content/reference_static_texts.md"
    content.parent.mkdir(parents=True)
    content.write_text("# Testi\n\n## appendix_placeholder\n\n", encoding="utf-8")

    expanded = auto_deck.expand_auto_slides(
        {
            "metadata": {"project_title": "Deliveroo vs Glovo", "subtitle": "Analisi", "authors": "Team", "date": "2026"},
            "auto_slides": {
                "enabled": True,
                "mode": "reference_order",
                "reference_texts": "slides/content/reference_static_texts.md",
                "task_count": 3,
                "appendix_assets_root": "slides/assets/appendices",
                "include_empty_appendices": True,
                "include_sources": False,
            },
            "slides": [],
        },
        {"cover", "section_divider_neutral", "text_only_neutral", "text_only_deliveroo", "text_only_glovo"},
    )

    titles = [slide.get("fields", {}).get("TEXT_TITLE") or slide.get("fields", {}).get("SECTION_NAME") for slide in expanded["slides"]]
    assert "Appendici" in titles
    assert "Appendice - Presentazione task user test - Le task" in titles
    assert "Deliveroo - Task 3" in titles
    assert "Glovo - Task 3" in titles
    assert "Appendice - Risultati questionario di usabilità" in titles


def test_resolve_template_id_variants() -> None:
    assert resolve_template_id("graph_full", "deliveroo") == "graph_full_deliveroo"
    assert resolve_template_id("graph_full", "glovo") == "graph_full_glovo"
    assert resolve_template_id("graph_full", "neutral") == "graph_full_neutral"
    assert resolve_template_id("comparison", "deliveroo") == "comparison_neutral"
    assert resolve_template_id("cover", "glovo") == "cover"
    assert resolve_template_id("text_only", "glovo") == "text_only_glovo"


def test_generate_slides_uses_theme_variants(tmp_path: Path) -> None:
    from PIL import Image
    from pptx import Presentation

    template = tmp_path / "template.pptx"
    config = tmp_path / "deck.yml"
    image = tmp_path / "graph.png"
    output = tmp_path / "out.pptx"

    prs = Presentation()
    for template_id, marker in [
        ("graph_full_deliveroo", "DELIVEROO_MARKER"),
        ("graph_full_glovo", "GLOVO_MARKER"),
        ("text_only_neutral", "NEUTRAL_MARKER"),
    ]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(100000, 100000, 3000000, 300000).text = f"TEMPLATE_ID:{template_id}"
        slide.shapes.add_textbox(100000, 500000, 5000000, 500000).text = marker
        slide.shapes.add_textbox(100000, 1000000, 5000000, 500000).text = "[GRAPH_TITLE][TEXT_TITLE]"
        slide.shapes.add_textbox(100000, 1600000, 3000000, 1200000).text = "[GRAPH_MAIN][TEXT_BODY]"
    prs.save(template)

    Image.new("RGB", (640, 360), "#00CCBC").save(image)
    config.write_text(
        f"""
deck:
  template: "{template.as_posix()}"
  output: "{output.as_posix()}"
slides:
  - template: graph_full
    theme: deliveroo
    fields:
      GRAPH_TITLE: "Deliveroo graph"
    images:
      GRAPH_MAIN: "{image.as_posix()}"
  - template: graph_full
    theme: glovo
    fields:
      GRAPH_TITLE: "Glovo graph"
    images:
      GRAPH_MAIN: "{image.as_posix()}"
  - template: text_only
    theme: neutral
    fields:
      TEXT_TITLE: "Neutral text"
      TEXT_BODY: "Body"
""",
        encoding="utf-8",
    )

    generate_slides(config, overwrite=True)
    generated = Presentation(output)
    text = "\n".join(
        shape.text
        for slide in generated.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "DELIVEROO_MARKER" in text
    assert "GLOVO_MARKER" in text
    assert "NEUTRAL_MARKER" in text


def test_generated_slide_text_uses_title_font_and_larger_section_titles(tmp_path: Path) -> None:
    from pptx import Presentation

    template = tmp_path / "template.pptx"
    config = tmp_path / "deck.yml"
    output = tmp_path / "out.pptx"

    prs = Presentation()
    section = prs.slides.add_slide(prs.slide_layouts[6])
    section.shapes.add_textbox(100000, 100000, 3000000, 300000).text = "TEMPLATE_ID:section_divider_neutral"
    section.shapes.add_textbox(100000, 500000, 5000000, 700000).text = "[SECTION_NAME]"

    text_slide = prs.slides.add_slide(prs.slide_layouts[6])
    text_slide.shapes.add_textbox(100000, 100000, 3000000, 300000).text = "TEMPLATE_ID:text_only_neutral"
    text_slide.shapes.add_textbox(100000, 600000, 5000000, 500000).text = "[TEXT_TITLE]"
    text_slide.shapes.add_textbox(100000, 1300000, 6000000, 2200000).text = "[TEXT_BODY]"
    prs.save(template)

    config.write_text(
        f"""
deck:
  template: "{template.as_posix()}"
  output: "{output.as_posix()}"
slides:
  - template_id: section_divider_neutral
    fields:
      SECTION_NAME: "Introduzione"
  - template_id: text_only_neutral
    fields:
      TEXT_TITLE: "Sintesi"
      TEXT_BODY: "Questo testo lungo usa lo stesso font dei titoli per valutare la resa complessiva."
""",
        encoding="utf-8",
    )

    generate_slides(config, overwrite=True)

    generated = Presentation(output)
    section_runs = [
        run
        for shape in generated.slides[0].shapes
        if getattr(shape, "has_text_frame", False) and "Introduzione" in shape.text
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    ]
    body_runs = [
        run
        for shape in generated.slides[1].shapes
        if getattr(shape, "has_text_frame", False) and "Questo testo lungo" in shape.text
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    ]

    assert section_runs
    assert body_runs
    assert all(run.font.name == DECK_FONT_FAMILY for run in body_runs)
    assert section_runs[0].font.name == DECK_FONT_FAMILY
    assert 18 <= section_runs[0].font.size.pt <= SECTION_TITLE_SIZE_SHORT


def test_generated_markdown_bullets_use_valid_unicode(tmp_path: Path) -> None:
    from pptx import Presentation

    template = tmp_path / "template.pptx"
    config = tmp_path / "deck.yml"
    output = tmp_path / "out.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(100000, 100000, 3000000, 300000).text = "TEMPLATE_ID:text_only_neutral"
    slide.shapes.add_textbox(100000, 600000, 5000000, 500000).text = "[TEXT_TITLE]"
    slide.shapes.add_textbox(100000, 1300000, 6000000, 2200000).text = "[TEXT_BODY]"
    prs.save(template)

    config.write_text(
        f"""
deck:
  template: "{template.as_posix()}"
  output: "{output.as_posix()}"
slides:
  - template_id: text_only_neutral
    fields:
      TEXT_TITLE: "Lista"
      TEXT_BODY: "Elementi:\\n- primo punto\\n- secondo punto"
""",
        encoding="utf-8",
    )

    generate_slides(config, overwrite=True)

    generated = Presentation(output)
    text = "\n".join(
        shape.text
        for shape in generated.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "• primo punto" in text
    assert "â€¢" not in text


def test_branded_section_titles_use_brand_text_color(tmp_path: Path) -> None:
    from pptx import Presentation

    template = tmp_path / "template.pptx"
    config = tmp_path / "deck.yml"
    output = tmp_path / "out.pptx"

    prs = Presentation()
    for template_id in ["section_divider_deliveroo", "section_divider_glovo"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(100000, 100000, 3000000, 300000).text = f"TEMPLATE_ID:{template_id}"
        slide.shapes.add_textbox(100000, 800000, 5000000, 800000).text = "[SECTION_NAME]"
    prs.save(template)

    config.write_text(
        f"""
deck:
  template: "{template.as_posix()}"
  output: "{output.as_posix()}"
slides:
  - template: section_divider
    theme: deliveroo
    fields:
      SECTION_NAME: "Deliveroo"
  - template: section_divider
    theme: glovo
    fields:
      SECTION_NAME: "Glovo"
""",
        encoding="utf-8",
    )

    generate_slides(config, overwrite=True)

    generated = Presentation(output)

    def title_color(slide_index: int, title: str) -> str:
        runs = [
            run
            for shape in generated.slides[slide_index].shapes
            if getattr(shape, "has_text_frame", False) and title in shape.text
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.text
        ]
        assert runs
        return str(runs[0].font.color.rgb)

    assert title_color(0, "Deliveroo") == DELIVEROO_COLOR
    assert title_color(1, "Glovo") == GLOVO_COLOR


def test_cover_title_uses_brand_colors_and_large_type(tmp_path: Path) -> None:
    from pptx import Presentation

    template = tmp_path / "template.pptx"
    config = tmp_path / "deck.yml"
    output = tmp_path / "out.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(100000, 100000, 3000000, 300000).text = "TEMPLATE_ID:cover"
    slide.shapes.add_textbox(100000, 700000, 8000000, 800000).text = "[PROJECT_TITLE]"
    prs.save(template)

    config.write_text(
        f"""
deck:
  template: "{template.as_posix()}"
  output: "{output.as_posix()}"
slides:
  - template_id: cover
    fields:
      PROJECT_TITLE: "Deliveroo vs Glovo"
""",
        encoding="utf-8",
    )

    generate_slides(config, overwrite=True)

    generated = Presentation(output)
    runs = [
        run
        for shape in generated.slides[0].shapes
        if getattr(shape, "has_text_frame", False) and "Deliveroo" in shape.text
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text
    ]
    assert [run.text for run in runs] == ["Deliveroo", " vs ", "Glovo"]
    assert runs[0].font.size.pt >= 60
    assert str(runs[0].font.color.rgb) == DELIVEROO_COLOR
    assert str(runs[2].font.color.rgb) == GLOVO_COLOR
