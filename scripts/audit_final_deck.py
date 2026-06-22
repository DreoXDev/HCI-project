from __future__ import annotations

import csv
import re
import sys
from collections import Counter
from pathlib import Path

import yaml
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "outputs" / "final" / "final_report.pptx"
REPORTS = ROOT / "outputs" / "reports"
CONFIG = ROOT / "config" / "slide_curation.yml"


def main() -> int:
    config = yaml.safe_load(CONFIG.read_text(encoding="utf-8")) if CONFIG.exists() else {}
    limits = config.get("limits", {})
    rows = audit_deck(PPTX, config)
    write_audit(rows)
    write_quality(rows, limits, config)
    print(REPORTS / "slide_audit.md")
    print(REPORTS / "deck_quality_report.md")
    return 0


def audit_deck(path: Path, config: dict) -> list[dict[str, object]]:
    prs = Presentation(str(path))
    title_counts: Counter[str] = Counter()
    raw_rows: list[dict[str, object]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = slide_title(slide)
        title_counts[title] += 1
        table_shapes = [shape for shape in slide.shapes if getattr(shape, "has_table", False)]
        row_count = max((len(shape.table.rows) - 1 for shape in table_shapes), default=0)
        col_count = max((len(shape.table.columns) for shape in table_shapes), default=0)
        text = "\n".join(getattr(shape, "text", "") for shape in slide.shapes if hasattr(shape, "text"))
        content_type = classify_content(slide, table_shapes, text)
        section = "appendix" if title.lower().startswith("appendice") else "main"
        raw_rows.append(
            {
                "slide": idx,
                "title": title,
                "section": section,
                "content_type": content_type,
                "table_rows": row_count,
                "table_columns": col_count,
                "too_many_rows": False,
                "too_many_columns": False,
                "likely_dataframe_dump": likely_dataframe_dump(row_count, col_count, text),
                "duplicate_candidate": False,
                "curation": curation_for(title, config),
            }
        )
    for row in raw_rows:
        row["duplicate_candidate"] = title_counts[str(row["title"])] > 1
        if row["section"] == "appendix":
            row["too_many_rows"] = int(row["table_rows"]) > int(config.get("limits", {}).get("appendix_max_rows", 10))
            row["too_many_columns"] = int(row["table_columns"]) > int(config.get("limits", {}).get("appendix_max_columns", 7))
        else:
            row["too_many_rows"] = int(row["table_rows"]) > int(config.get("limits", {}).get("main_max_rows", 6))
            row["too_many_columns"] = int(row["table_columns"]) > int(config.get("limits", {}).get("main_max_columns", 6))
    return raw_rows


def slide_title(slide) -> str:
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = " ".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
        if text:
            candidates.append((int(getattr(shape, "top", 0)), text))
    if not candidates:
        return "(senza titolo)"
    return sorted(candidates, key=lambda item: item[0])[0][1]


def classify_content(slide, tables: list, text: str) -> str:
    pictures = [shape for shape in slide.shapes if getattr(shape, "shape_type", None) == 13]
    if tables and pictures:
        return "grafico+tabella"
    if tables:
        return "tabella"
    if pictures:
        return "grafico"
    return "testo"


def likely_dataframe_dump(rows: int, cols: int, text: str) -> bool:
    technical_tokens = ["p_value", "effect_size", "ci95", "participant_id", "transformed", "missing_count"]
    return rows > 10 or cols > 7 or any(token in text for token in technical_tokens)


def curation_for(title: str, config: dict) -> str:
    for key in ["main_keep", "main_refactor", "appendix_keep", "appendix_refactor", "delete_duplicate", "export_only"]:
        if any(item in title for item in config.get(key, []) if isinstance(item, str)):
            return key
    if title.lower().startswith("appendice"):
        return "appendix_keep"
    return "main_keep"


def write_audit(rows: list[dict[str, object]]) -> None:
    REPORTS.mkdir(parents=True, exist_ok=True)
    csv_path = REPORTS / "slide_audit.csv"
    with csv_path.open("w", newline="", encoding="utf-8-sig") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0].keys()) if rows else [])
        writer.writeheader()
        writer.writerows(rows)
    lines = ["# Slide audit", "", f"- Slide analizzate: {len(rows)}", ""]
    flagged = [row for row in rows if row["too_many_rows"] or row["too_many_columns"] or row["likely_dataframe_dump"]]
    lines.append(f"- Slide con warning leggibilita: {len(flagged)}")
    lines.extend(["", "## Warning principali", ""])
    for row in flagged[:40]:
        lines.append(
            f"- Slide {row['slide']}: {row['title']} "
            f"(righe={row['table_rows']}, colonne={row['table_columns']}, dump={row['likely_dataframe_dump']})"
        )
    (REPORTS / "slide_audit.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_quality(rows: list[dict[str, object]], limits: dict, config: dict) -> None:
    max_slides = int(limits.get("deck_max_slides_warning", 160))
    appendix = [row for row in rows if row["section"] == "appendix"]
    warnings = []
    if len(rows) > max_slides:
        warnings.append(f"Deck sopra target: {len(rows)} slide > {max_slides}.")
    for row in rows:
        if is_density_exempt(str(row["title"]), config):
            continue
        if row["section"] == "main" and (row["too_many_rows"] or row["too_many_columns"]):
            warnings.append(f"Slide principale troppo densa: {row['slide']} - {row['title']}.")
        if row["section"] == "appendix" and (row["too_many_rows"] or row["too_many_columns"]):
            warnings.append(f"Appendice troppo densa: {row['slide']} - {row['title']}.")
    lines = [
        "# Deck quality report",
        "",
        f"- Slide totali: {len(rows)}",
        f"- Slide appendice: {len(appendix)}",
        f"- Target massimo warning: {max_slides}",
        f"- Esito: {'WARNING' if warnings else 'OK'}",
        "",
        "## Warning",
        "",
    ]
    lines.extend([f"- {warning}" for warning in warnings] or ["- Nessuno."])
    (REPORTS / "deck_quality_report.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def is_density_exempt(title: str, config: dict) -> bool:
    return any(prefix in title for prefix in config.get("density_exempt_titles", []))


if __name__ == "__main__":
    sys.exit(main())
