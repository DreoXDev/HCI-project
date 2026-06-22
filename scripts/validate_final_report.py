from __future__ import annotations

import math
import re
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
PLACEHOLDER_TOKENS = ["INSERIRE", "TODO", "PLACEHOLDER", "screenshot qui", "SCREENSHOT QUI"]
REQUIRED_FILES = [
    "outputs/final/final_report.pptx",
    "outputs/final/final_report.pdf",
    "data/user_testing_times.csv",
    "data/user_profiles.csv",
    "data/processed/user_task_trials_normalized.csv",
    "data/processed/final_report/questionnaire_item_descriptive_stats.csv",
    "data/processed/final_report/ueq_item_mapping.csv",
    "data/processed/final_report/ueq_scoring_method.md",
    "data/processed/final_report/ueq_scale_validation.md",
    "data/processed/final_report/ueq_scale_summary.csv",
    "outputs/tables/user_profile_scoring_method.md",
]
REQUIRED_NONEMPTY_TABLES = [
    "data/user_testing_times.csv",
    "data/user_profiles.csv",
    "data/processed/user_task_trials_normalized.csv",
    "outputs/tables/user_task_trials_full.csv",
    "outputs/tables/user_profiles_slide.csv",
    "outputs/tables/user_test_efficiency_by_task.csv",
    "outputs/tables/user_test_efficiency_comparison.csv",
    "outputs/tables/user_test_efficiency_comparison_slide.csv",
    "outputs/tables/user_test_effectiveness_by_task.csv",
    "outputs/tables/user_test_assistance_errors.csv",
    "outputs/tables/user_test_assistance_events_slide.csv",
    "outputs/tables/user_testing_by_expertise_group.csv",
    "data/processed/final_report/user_test_times_unified.csv",
    "data/processed/final_report/user_test_times_summary.csv",
    "data/processed/final_report/task_efficiency_stats.csv",
    "data/processed/final_report/questionnaire_item_descriptive_stats.csv",
    "data/processed/final_report/ueq_item_mapping.csv",
    "data/processed/final_report/ueq_item_summary.csv",
    "data/processed/final_report/ueq_scale_summary.csv",
]
REQUIRED_UEQ_ASSETS = [
    "slides/assets/generated/ueq/ueq_distribution_deliveroo.png",
    "slides/assets/generated/ueq/ueq_distribution_glovo.png",
    "slides/assets/generated/ueq/ueq_mean_results_deliveroo.png",
    "slides/assets/generated/ueq/ueq_mean_results_glovo.png",
    "slides/assets/generated/ueq/ueq_benchmark_deliveroo.png",
    "slides/assets/generated/ueq/ueq_benchmark_glovo.png",
    "slides/assets/generated/ueq/ueq_scale_comparison_deliveroo_vs_glovo.png",
]
REQUIRED_USER_TEST_ASSETS = [
    "outputs/charts/user_demographics_age_gender.png",
    "outputs/charts/users_age_pie.png",
    "outputs/charts/users_gender_pie.png",
    "outputs/charts/users_occupation_pie.png",
    "outputs/charts/users_delivery_familiarity_pie.png",
    "outputs/charts/experts_age_pie.png",
    "outputs/charts/experts_gender_pie.png",
    "outputs/charts/experts_occupation_pie.png",
    "outputs/charts/experts_delivery_familiarity_pie.png",
    "outputs/charts/user_profile_delivery_familiarity.png",
    "outputs/charts/user_test_effectiveness_strict.png",
    "outputs/charts/user_test_effectiveness_extended.png",
    "outputs/charts/user_test_assisted_tasks.png",
    "outputs/charts/user_test_efficiency_task_1.png",
    "outputs/charts/user_test_efficiency_task_2.png",
    "outputs/charts/user_test_efficiency_task_3.png",
    "outputs/charts/user_test_efficiency_overall.png",
    "outputs/charts/effectiveness_outcome_matrix.png",
    "outputs/charts/efficiency_summary.png",
    "outputs/charts/efficiency_task_1_boxplot.png",
    "outputs/charts/efficiency_task_1_paired_lines.png",
    "outputs/charts/efficiency_task_2_boxplot.png",
    "outputs/charts/efficiency_task_2_paired_lines.png",
    "outputs/charts/efficiency_task_3_boxplot.png",
    "outputs/charts/efficiency_task_3_paired_lines.png",
    "outputs/charts/ueq_benchmark_comparison.png",
    "outputs/charts/subgroup_ueq_heatmap.png",
]
REQUIRED_SLIDE_TITLES = [
    "Profilo degli utenti coinvolti",
    "Composizione del campione utenti",
    "Efficacia - matrice esiti utenti/task",
    "Efficacia relativa - sintesi",
    "Efficacia assoluta - sintesi",
    "Efficienza - riepilogo",
    "Efficienza statistica - sintesi",
    "UEQ D9",
    "UEQ benchmark - confronto sintetico",
    "Confronto statistico complessivo",
    "Task assistite e issue",
]


def main() -> int:
    errors: list[str] = []
    for rel in [*REQUIRED_FILES, *REQUIRED_UEQ_ASSETS, *REQUIRED_USER_TEST_ASSETS]:
        path = ROOT / rel
        if not path.exists() or path.stat().st_size == 0:
            errors.append(f"Missing required file: {rel}")

    pptx = ROOT / "outputs/final/final_report.pptx"
    if pptx.exists():
        errors.extend(_validate_pptx_text(pptx))
        errors.extend(_validate_required_slide_titles(pptx))

    for rel in REQUIRED_NONEMPTY_TABLES:
        errors.extend(_validate_nonempty_table(ROOT / rel, rel))

    trials = ROOT / "data/processed/user_task_trials_normalized.csv"
    if trials.exists():
        errors.extend(_validate_trials(trials))
    real_times = ROOT / "data/user_testing_times.csv"
    profiles = ROOT / "data/user_profiles.csv"
    if real_times.exists():
        errors.extend(_validate_real_times(real_times))
    if profiles.exists() and real_times.exists():
        errors.extend(_validate_profiles(profiles, real_times))

    descriptives = ROOT / "data/processed/final_report/questionnaire_item_descriptive_stats.csv"
    if descriptives.exists():
        df = pd.read_csv(descriptives, encoding="utf-8-sig")
        if df.get("item_id", pd.Series(dtype=float)).nunique() < 26:
            errors.append("Questionnaire descriptives cover fewer than 26 items.")
        if len(df) < 52:
            errors.append("Questionnaire descriptives have fewer than 52 rows.")

    _write_review_report(errors)
    if errors:
        print("STATUS: FAIL")
        for error in errors:
            print(f"- {error}")
        return 1
    print("STATUS: OK")
    return 0


def _validate_pptx_text(path: Path) -> list[str]:
    errors: list[str] = []
    prs = Presentation(str(path))
    for idx, slide in enumerate(prs.slides, start=1):
        text = "\n".join(
            "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs)
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        lowered = text.lower()
        for token in PLACEHOLDER_TOKENS:
            if token.upper() == "TODO":
                matched = re.search(r"\bTODO\b", text, flags=re.IGNORECASE)
            else:
                matched = token.lower() in lowered
            if matched:
                errors.append(f"Slide {idx} contains placeholder token: {token}")
        for token in ["nan", "none", "inf"]:
            if f" {token} " in f" {lowered} ":
                errors.append(f"Slide {idx} contains visible invalid value: {token}")
    return errors


def _validate_required_slide_titles(path: Path) -> list[str]:
    prs = Presentation(str(path))
    titles: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = " ".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
                if text:
                    titles.append(text)
                    break
    return [f"Missing required slide title: {title}" for title in REQUIRED_SLIDE_TITLES if not any(title in existing for existing in titles)]


def _validate_nonempty_table(path: Path, rel: str) -> list[str]:
    if not path.exists():
        return [f"Missing table: {rel}"]
    try:
        df = pd.read_csv(path, encoding="utf-8-sig")
    except Exception as exc:
        return [f"Cannot read table {rel}: {exc}"]
    if df.empty:
        return [f"Table has header only or zero rows: {rel}"]
    if df.replace("", pd.NA).dropna(how="all").empty:
        return [f"Table has no non-empty data rows: {rel}"]
    return []


def _validate_trials(path: Path) -> list[str]:
    errors: list[str] = []
    df = pd.read_csv(path, encoding="utf-8-sig")
    required = {"participant_id", "app", "task_id", "time_seconds", "outcome", "completed", "correct", "assisted"}
    missing = sorted(required - set(df.columns))
    if missing:
        errors.append(f"Trial table missing columns: {', '.join(missing)}")
        return errors
    if len(df) < 144:
        errors.append(f"Trial table has fewer than 144 rows: {len(df)}")
    valid_outcomes = {"success", "assisted_success", "partial_success", "failure", "timeout"}
    if not df["outcome"].isin(valid_outcomes).any():
        errors.append("Trial table has no valid outcome values.")
    if pd.to_numeric(df["time_seconds"], errors="coerce").dropna().empty:
        errors.append("Trial table has no numeric time_seconds.")
    pivot = df[df["outcome"].isin(valid_outcomes)].pivot_table(index=["participant_id", "task_id"], columns="app", values="time_seconds", aggfunc="first")
    if not {"Deliveroo", "Glovo"}.issubset(set(map(str, pivot.columns))):
        errors.append("No paired Deliveroo/Glovo task times available.")
    else:
        pairs = pivot[["Deliveroo", "Glovo"]].apply(pd.to_numeric, errors="coerce").dropna()
        if len(pairs) < 24:
            errors.append(f"Paired task rows insufficient for statistical tests: {len(pairs)}")
    for column in ["time_seconds"]:
        numeric = pd.to_numeric(df[column], errors="coerce")
        if numeric.map(lambda value: pd.notna(value) and not math.isfinite(float(value))).any():
            errors.append(f"Trial table contains non-finite values in {column}.")
    return errors


def _validate_real_times(path: Path) -> list[str]:
    errors: list[str] = []
    df = pd.read_csv(path, encoding="utf-8-sig")
    if len(df) != 144:
        errors.append(f"data/user_testing_times.csv must contain 144 rows, found {len(df)}.")
    key = ["user_id", "app", "task"]
    if df.duplicated(key).any():
        errors.append("Duplicate user_id/app/task combinations in user_testing_times.")
    expected_users = {f"U{i}" for i in range(1, 25)}
    expected_apps = {"Deliveroo", "Glovo"}
    expected_tasks = {1, 2, 3}
    actual = {(row.user_id, row.app, int(row.task)) for row in df.itertuples(index=False)}
    missing = [(u, a, t) for u in expected_users for a in expected_apps for t in expected_tasks if (u, a, t) not in actual]
    if missing:
        errors.append(f"Missing user/app/task combinations: {len(missing)}")
    if pd.to_numeric(df["time_seconds"], errors="coerce").isna().any():
        errors.append("Some time_seconds values are not numeric.")
    sample = df[df["time_raw"] == "01:41.07"]
    if sample.empty or abs(float(sample.iloc[0]["time_seconds"]) - 101.07) > 0.001:
        errors.append("Time conversion check failed for 01:41.07 -> 101.07.")
    assisted = df[(df["assistance"].astype(str) == "verbal_help") | (df["outcome"].astype(str) == "assisted_success")]
    if len(assisted) < 11:
        errors.append(f"Expected at least 11 assisted tasks, found {len(assisted)}.")
    expected_assisted = {
        ("U1", "Deliveroo", 3),
        ("U1", "Glovo", 3),
        ("U2", "Deliveroo", 3),
        ("U3", "Deliveroo", 2),
        ("U5", "Deliveroo", 3),
        ("U5", "Glovo", 3),
        ("U14", "Deliveroo", 3),
        ("U15", "Deliveroo", 3),
        ("U16", "Deliveroo", 3),
        ("U22", "Deliveroo", 3),
        ("U24", "Deliveroo", 3),
    }
    actual_assisted = {(row.user_id, row.app, int(row.task)) for row in assisted.itertuples(index=False)}
    missing_assisted = sorted(expected_assisted - actual_assisted)
    if missing_assisted:
        errors.append(f"Missing assisted task codings: {missing_assisted}")
    u5_glovo = df[(df["user_id"] == "U5") & (df["app"] == "Glovo") & (df["task"].astype(int) == 3)]
    if u5_glovo.empty or str(u5_glovo.iloc[0]["outcome"]) != "assisted_success":
        errors.append("U5 Glovo task 3 must be coded as assisted_success.")
    return errors


def _validate_profiles(profile_path: Path, times_path: Path) -> list[str]:
    errors: list[str] = []
    profiles = pd.read_csv(profile_path, encoding="utf-8-sig")
    times = pd.read_csv(times_path, encoding="utf-8-sig")
    if len(profiles) != 24:
        errors.append(f"user_profiles.csv must contain 24 users, found {len(profiles)}.")
    profile_ids = set(profiles["user_id"].astype(str))
    time_ids = set(times["user_id"].astype(str))
    if profile_ids != time_ids:
        errors.append("User IDs in user_profiles.csv do not match user_testing_times.csv.")
    if pd.to_numeric(profiles["delivery_familiarity"], errors="coerce").dropna().shape[0] < 24:
        errors.append("delivery_familiarity is incomplete in user_profiles.csv.")
    for column in ["delivery_familiarity", "digital_familiarity"]:
        if column not in profiles.columns:
            continue
        numeric = pd.to_numeric(profiles[column], errors="coerce")
        if numeric.dropna().min() < 1 or numeric.dropna().max() > 3:
            errors.append(f"{column} must remain on original 1-3 scale.")
    return errors


def _write_review_report(errors: list[str]) -> None:
    report = ROOT / "outputs/reports/final_report_review_after_reorder.md"
    report.parent.mkdir(parents=True, exist_ok=True)
    pptx = ROOT / "outputs/final/final_report.pptx"
    titles: list[str] = []
    if pptx.exists():
        prs = Presentation(str(pptx))
        for slide in prs.slides:
            title = ""
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = " ".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
                    if text:
                        title = text
                        break
            titles.append(title or "(senza titolo)")
    lines = [
        "# Review finale dopo reorder",
        "",
        f"- Slide generate: {len(titles)}",
        "- Slide rimosse/ridotte: duplicati demografici utenti nella sezione Questionario; separatore autonomo Dark pattern.",
        "- Slide riordinate: sintesi euristica spostata dopo dark pattern; metodo UEQ anticipato; campione utenti collocato nel blocco Test utente.",
        "- Grafici sostituiti: pie chart demografiche per valutatori e utenti; matrice utenti su scala originale 1-3.",
        "- Tabelle formattate: confronto statistico task compatto; tabella eventi assistenze e issue.",
        f"- Warning rimasti: {len(errors)}",
        "",
        "## Titoli slide",
        "",
    ]
    lines.extend(f"{idx}. {title}" for idx, title in enumerate(titles, start=1))
    if errors:
        lines.extend(["", "## Errori quality gate", ""])
        lines.extend(f"- {error}" for error in errors)
    report.write_text("\n".join(lines) + "\n", encoding="utf-8")


if __name__ == "__main__":
    sys.exit(main())
