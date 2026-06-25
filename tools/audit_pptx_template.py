from __future__ import annotations

import argparse
import json
import re
import sys
from collections import Counter
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from src.config import resolve_path
from src.slide_export.auto_deck import expand_auto_slides
from src.slide_export.pptx_generator import (
    _expand_paginated_table_specs,
    _load_yaml,
    _resolved_spec_template_id,
    _shape_text,
    _template_slides_by_id,
)
from src.slide_export.template_variants import REQUIRED_SHAPES, REQUIRED_TEMPLATE_IDS, base_template_for


OPTIONAL_BASES = {"findings", "task_results", "ueq_question", "sources", "final_verdict"}
OPTIONAL_TEMPLATE_IDS = {
    "section_divider_deliveroo",
    "section_divider_glovo",
}


def _shape_name(shape: Any) -> str:
    return str(getattr(shape, "name", "") or "")


def _extract_template_ids(text: str) -> list[str]:
    return re.findall(r"TEMPLATE_ID\s*:\s*([A-Za-z0-9_-]+)", text)


def _extract_placeholders(text: str) -> list[str]:
    bracketed = re.findall(r"\[([A-Z0-9_ /.-]+)\]", text)
    bare = re.findall(r"\b([A-Z][A-Z0-9_]{2,})\b", text)
    ignored = {"TEMPLATE_ID"}
    return sorted({item.strip() for item in [*bracketed, *bare] if item.strip() and item.strip() not in ignored})


def _slide_inventory(presentation: Any) -> list[dict[str, Any]]:
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        shape_rows = []
        all_text = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            all_text.append(text)
            shape_rows.append(
                {
                    "name": _shape_name(shape),
                    "type": int(getattr(shape, "shape_type", -1)),
                    "text": text.strip()[:200],
                    "placeholders": _extract_placeholders(text),
                    "is_picture": int(getattr(shape, "shape_type", -1)) == 13,
                }
            )
        joined = "\n".join(all_text)
        ids = _extract_template_ids(joined)
        slides.append(
            {
                "slide_index": index,
                "template_ids": ids,
                "template_id": ids[0] if ids else None,
                "shape_count": len(shape_rows),
                "picture_count": sum(1 for row in shape_rows if row["is_picture"]),
                "shape_names": [row["name"] for row in shape_rows if row["name"]],
                "text_placeholders": sorted({ph for row in shape_rows for ph in row["placeholders"]}),
                "shapes": shape_rows,
            }
        )
    return slides


def _used_template_counts(config_path: Path, template_ids: set[str]) -> Counter[str]:
    deck_config = _load_yaml(config_path)
    expanded = expand_auto_slides(deck_config, template_ids)
    requested = _expand_paginated_table_specs(expanded.get("slides", []))
    counts: Counter[str] = Counter()
    for spec in requested:
        if spec.get("template_id") == "full_slide_image":
            counts["full_slide_image"] += 1
        else:
            counts[_resolved_spec_template_id(spec)] += 1
    return counts


def _classify(template_id: str, used_count: int) -> tuple[str, str]:
    base = base_template_for(template_id)
    required_placeholders = REQUIRED_SHAPES.get(base, [])
    if used_count:
        return "USED_REQUIRED", "keep"
    if template_id in OPTIONAL_TEMPLATE_IDS or base in OPTIONAL_BASES:
        return "USED_OPTIONAL", "keep in main or move to legacy after config cleanup"
    if template_id in REQUIRED_TEMPLATE_IDS:
        return "LEGACY_COMPAT", "keep until code support matrix is narrowed"
    return "UNUSED_REMOVE", "move out of active template"


def build_audit(template_path: Path, config_path: Path) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(str(template_path))
    template_map = _template_slides_by_id(presentation)
    inventory = _slide_inventory(presentation)
    used_counts = _used_template_counts(config_path, set(template_map))

    compatibility = []
    for row in inventory:
        template_id = row["template_id"]
        if not template_id:
            compatibility.append({**row, "status": "UNUSED_REMOVE", "action": "remove or mark as guide/non-generating", "used_count": 0, "required_placeholders": []})
            continue
        status, action = _classify(template_id, used_counts.get(template_id, 0))
        compatibility.append(
            {
                **row,
                "used_count": used_counts.get(template_id, 0),
                "required_placeholders": REQUIRED_SHAPES.get(base_template_for(template_id), []),
                "status": status,
                "action": action,
            }
        )
    return {
        "template_path": str(template_path),
        "config_path": str(config_path),
        "slide_count": len(inventory),
        "template_id_counts": dict(Counter(row["template_id"] or "<missing>" for row in inventory)),
        "template_ids": sorted(template_map),
        "used_template_counts": dict(sorted(used_counts.items())),
        "compatibility": compatibility,
    }


def _markdown_table(headers: list[str], rows: list[list[Any]]) -> list[str]:
    out = ["| " + " | ".join(headers) + " |", "| " + " | ".join(["---"] * len(headers)) + " |"]
    for row in rows:
        out.append("| " + " | ".join(str(cell).replace("\n", " ") for cell in row) + " |")
    return out


def write_markdown(audit: dict[str, Any], path: Path) -> None:
    lines = [
        "# PowerPoint Template Audit",
        "",
        f"- Template: `{audit['template_path']}`",
        f"- Slide config: `{audit['config_path']}`",
        f"- Template slides: {audit['slide_count']}",
        f"- TEMPLATE_ID found: {len(audit['template_ids'])}",
        "",
        "## TEMPLATE_ID Inventory",
        "",
        *_markdown_table(["TEMPLATE_ID", "count"], [[key, value] for key, value in sorted(audit["template_id_counts"].items())]),
        "",
        "## Compatibility Matrix",
        "",
        *_markdown_table(
            ["slide", "TEMPLATE_ID", "used", "required placeholders", "text placeholders found", "status", "action"],
            [
                [
                    row["slide_index"],
                    row["template_id"] or "-",
                    row["used_count"],
                    ", ".join(row["required_placeholders"]),
                    ", ".join(row["text_placeholders"]),
                    row["status"],
                    row["action"],
                ]
                for row in audit["compatibility"]
            ],
        ),
        "",
        "## Cleanup Recommendation",
        "",
        "The final-report pipeline currently uses only layouts marked `USED_REQUIRED`. Layouts marked `USED_OPTIONAL` are retained because optional auto-generation paths can still use them when enabled. Move them to the legacy template only after disabling or documenting those optional paths.",
        "",
        "## Student-Safe Editing Rules",
        "",
        "- Safe to edit: colors, backgrounds, typography, decorative shapes and image styling.",
        "- Do not rename or delete: `TEMPLATE_ID` markers and placeholder shape names/text referenced by the generator.",
        "- After editing: run `python -m src.cli validate-template` and regenerate the deck.",
    ]
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> None:
    parser = argparse.ArgumentParser(description="Audit PowerPoint template TEMPLATE_IDs and placeholders.")
    parser.add_argument("--template", default="slides/templates/Deliveroo_vs_Glovo_clean_python_ready_template.pptx")
    parser.add_argument("--config", default="slides/config/slide_deck.yml")
    parser.add_argument("--json-out", default="outputs/template_audit.json")
    parser.add_argument("--md-out", default="docs/template-audit.md")
    args = parser.parse_args()

    template_path = resolve_path(args.template)
    config_path = resolve_path(args.config)
    audit = build_audit(template_path, config_path)
    json_path = resolve_path(args.json_out)
    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps(audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    write_markdown(audit, resolve_path(args.md_out))
    print(resolve_path(args.md_out))
    print(json_path)


if __name__ == "__main__":
    main()
