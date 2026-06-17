from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from ..config import resolve_path


BANNED_PATTERNS = [
    "Placeholder",
    "Spazio per",
    "Output generato",
    "Verificare la lettura",
    "TODO",
    "FIXME",
    "Lorem ipsum",
    "outputs/",
    "slides/assets/",
    "data/raw/",
    "INSERIRE QUI",
    "Asset generati",
    "[object Object]",
    "NaN",
    "None",
    "null",
    "undefined",
    "PARTIAL_DATA",
    "rimuovendo il menu appena aggiunto",
]
TOKEN_PATTERNS = {"TODO", "FIXME", "NaN", "None", "null", "undefined"}


@dataclass
class DeckTextFinding:
    slide_number: int
    pattern: str
    text: str


def audit_final_deck_text(
    pptx_path: str | Path = "outputs/slides/final_report.pptx",
    output_path: str | Path = "outputs/quality/final_report_text_audit.md",
) -> list[DeckTextFinding]:
    from pptx import Presentation

    target = resolve_path(pptx_path)
    findings: list[DeckTextFinding] = []
    if not target.exists():
        findings.append(DeckTextFinding(0, "missing_file", f"PPTX non trovato: {target}"))
        _write_audit(findings, output_path)
        return findings
    presentation = Presentation(str(target))
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for text in _slide_texts(slide):
            normalized = text.replace("\\", "/")
            for pattern in BANNED_PATTERNS:
                if _contains_banned(normalized, pattern):
                    findings.append(DeckTextFinding(slide_number, pattern, _compact(normalized)))
            if "`" in normalized:
                findings.append(DeckTextFinding(slide_number, "backtick Markdown", _compact(normalized)))
    _write_audit(findings, output_path)
    return findings


def _slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        texts.append(text)
    return texts


def _write_audit(findings: list[DeckTextFinding], output_path: str | Path) -> Path:
    target = resolve_path(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    lines = ["# Final Report Text Audit", ""]
    if not findings:
        lines.append("STATUS: OK")
    else:
        lines.extend(["STATUS: FAIL", "", "| slide | pattern | text |", "|---:|---|---|"])
        for finding in findings:
            safe_text = finding.text.replace("|", "\\|")
            lines.append(f"| {finding.slide_number} | {finding.pattern} | {safe_text} |")
    target.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return target


def _contains_banned(text: str, pattern: str) -> bool:
    if pattern in TOKEN_PATTERNS:
        return re.search(rf"(?<![A-Za-z0-9_]){re.escape(pattern)}(?![A-Za-z0-9_])", text, flags=re.IGNORECASE) is not None
    return pattern.casefold() in text.casefold()


def _compact(text: str, limit: int = 180) -> str:
    value = re.sub(r"\s+", " ", str(text)).strip()
    return value if len(value) <= limit else value[:limit].rstrip()
