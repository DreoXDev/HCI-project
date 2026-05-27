from __future__ import annotations

import csv
import re
import shutil
from copy import deepcopy
from dataclasses import dataclass, field
from datetime import datetime
from math import ceil
from pathlib import Path
from typing import Any

import yaml

from ..config import resolve_path
from ..text_generation.italian import italian_display_text
from .auto_deck import expand_auto_slides
from .tables import read_display_table, table_specs_from_paginated_table
from .template_variants import REQUIRED_SHAPES, REQUIRED_TEMPLATE_IDS, base_template_for, resolve_template_id

DECK_FONT_FAMILY = "Sora"
DELIVEROO_COLOR = "00CCBC"
GLOVO_COLOR = "FFC244"
SECTION_TITLE_SIZE_SHORT = 68
SECTION_TITLE_SIZE_MEDIUM = 60
SECTION_TITLE_SIZE_LONG = 52
COVER_TITLE_SIZE = 66


class SlideGenerationError(RuntimeError):
    """Raised when the PPTX deck cannot be generated with actionable context."""


@dataclass
class SlideGenerationResult:
    template: Path
    config: Path
    output: Path
    generated_slides: int
    missing_optional_assets: list[Path] = field(default_factory=list)


def generate_slides(
    config_path: str | Path = "slides/config/slide_deck.yml",
    *,
    template_path: str | Path | None = None,
    output_path: str | Path | None = None,
    overwrite: bool = False,
    timestamp: bool = False,
) -> SlideGenerationResult:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise SlideGenerationError(
            "Missing dependency: python-pptx.\n\nFix:\nRun `python -m pip install -r requirements.txt`."
        ) from exc

    config_file = resolve_path(config_path)
    if not config_file.exists():
        raise SlideGenerationError(f"Slide config not found:\n{config_file}")

    deck_config = _load_yaml(config_file)
    deck = deck_config.get("deck", {})
    template = resolve_path(template_path or deck.get("template") or "slides/templates/master_template.pptx")
    output = resolve_path(output_path or deck.get("output") or "outputs/slides/final_report.pptx")
    if timestamp:
        output = output.with_name(f"{output.stem}_{datetime.now().strftime('%Y-%m-%d_%H-%M')}{output.suffix}")

    _validate_paths(template, output, overwrite)
    presentation = Presentation(str(template))
    template_map = _template_slides_by_id(presentation)
    deck_config = expand_auto_slides(deck_config, set(template_map))
    requested = deck_config.get("slides", [])
    requested = _expand_paginated_table_specs(requested)
    _validate_requested_slides(requested, template_map)
    _validate_assets(deck_config)

    original_slide_count = len(presentation.slides)
    missing_optional: list[Path] = []
    for spec in requested:
        template_id = _resolved_spec_template_id(spec)
        source_slide = template_map[template_id]
        slide = _duplicate_slide(presentation, source_slide)
        _remove_template_metadata_shapes(slide)
        fields = _collect_fields(spec)
        _replace_text_fields(slide, fields)
        _replace_named_non_text_fields_with_text(slide, fields, set((spec.get("images") or {}).keys()))
        for placeholder, image in (spec.get("images") or {}).items():
            image_path = resolve_path(image)
            if image_path.exists():
                replace_placeholder_with_image(slide, placeholder, image_path)
            else:
                missing_optional.append(image_path)
        text_box = spec.get("text_box")
        if text_box:
            replace_placeholder_with_text_box(slide, text_box["placeholder"], str(text_box["text"]))
        table = spec.get("table")
        if table:
            replace_placeholder_with_table(
                slide,
                table["placeholder"],
                resolve_path(table["source"]),
                start_row=int(table.get("start_row") or 0),
                max_rows=int(table["max_rows"]) if table.get("max_rows") else None,
                font_size=float(table.get("font_size") or 7.5),
                header_font_size=float(table.get("header_font_size") or table.get("font_size") or 7.5),
                max_cell_chars=int(table.get("max_cell_chars") or 70),
                column_widths=table.get("column_widths"),
            )
        _clear_unresolved_placeholders(slide)
        _style_and_fit_slide_text_dark(slide, template_id=template_id)

    _remove_template_slides(presentation, original_slide_count)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output))
    return SlideGenerationResult(
        template=template,
        config=config_file,
        output=output,
        generated_slides=len(requested),
        missing_optional_assets=missing_optional,
    )


def validate_slide_assets(
    config_path: str | Path = "slides/config/slide_deck.yml",
    *,
    template_path: str | Path | None = None,
) -> list[str]:
    config_file = resolve_path(config_path)
    if not config_file.exists():
        return [f"Slide config not found: {config_file}"]
    try:
        deck_config = _load_yaml(config_file)
        template = resolve_path(template_path or deck_config.get("deck", {}).get("template") or "slides/templates/master_template.pptx")
        messages = []
        if not template.exists():
            messages.append(f"Template not found: {template}")
        else:
            try:
                from pptx import Presentation

                template_map = _template_slides_by_id(Presentation(str(template)))
                deck_config = expand_auto_slides(deck_config, set(template_map))
            except ModuleNotFoundError:
                pass
        messages.extend(_asset_errors(deck_config))
        return messages
    except Exception as exc:
        return [f"Invalid slide config: {exc}"]


def replace_placeholder_with_image(slide: Any, placeholder_name: str, image_path: str | Path) -> None:
    from PIL import Image

    shape = _find_placeholder_shape(slide, placeholder_name)
    if shape is None:
        shape = _find_visual_image_placeholder(slide, placeholder_name)
    if shape is None:
        left, top, width, height = _fallback_image_bounds(placeholder_name)
    else:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    rounded_path = _rounded_image_asset(image_path)
    with Image.open(rounded_path) as image:
        img_width, img_height = image.size
    box_ratio = width / height
    image_ratio = img_width / img_height
    if image_ratio > box_ratio:
        new_width = width
        new_height = int(width / image_ratio)
    else:
        new_height = height
        new_width = int(height * image_ratio)
    if shape is not None:
        _remove_shape(shape)
    slide.shapes.add_picture(
        str(rounded_path),
        left + int((width - new_width) / 2),
        top + int((height - new_height) / 2),
        width=new_width,
        height=new_height,
    )


def _find_visual_image_placeholder(slide: Any, placeholder_name: str) -> Any | None:
    normalized = _normalize_placeholder(placeholder_name)
    pictures = [
        shape
        for shape in slide.shapes
        if shape.shape_type == 13
        and int(getattr(shape, "top", 0)) > 900000
        and int(getattr(shape, "width", 0)) > 1200000
        and int(getattr(shape, "height", 0)) > 900000
    ]
    if not pictures:
        return None
    if "LEFT" in normalized:
        return sorted(pictures, key=lambda shape: int(shape.left))[0]
    if "RIGHT" in normalized:
        return sorted(pictures, key=lambda shape: int(shape.left))[-1]
    if normalized == "GRAPHMAIN":
        return max(pictures, key=lambda shape: int(shape.width) * int(shape.height))
    if "MINI" in normalized or "SUMMARY" in normalized:
        return min(pictures, key=lambda shape: int(shape.width) * int(shape.height))
    return None


def _rounded_image_asset(image_path: str | Path, radius: int = 34) -> Path:
    from PIL import Image, ImageDraw

    source = resolve_path(image_path)
    target_dir = resolve_path("outputs/slide_assets/rounded")
    target_dir.mkdir(parents=True, exist_ok=True)
    target = target_dir / f"{source.stem}_rounded.png"
    if target.exists() and target.stat().st_mtime >= source.stat().st_mtime:
        return target
    with Image.open(source).convert("RGBA") as image:
        mask = Image.new("L", image.size, 0)
        draw = ImageDraw.Draw(mask)
        draw.rounded_rectangle((0, 0, image.width, image.height), radius=radius, fill=255)
        image.putalpha(mask)
        image.save(target)
    return target


def _fallback_image_bounds(placeholder_name: str) -> tuple[int, int, int, int]:
    normalized = _normalize_placeholder(placeholder_name)
    if normalized == "GRAPHMAIN":
        return 650000, 2180000, 10900000, 4250000
    if "LEFT" in normalized:
        return 650000, 1600000, 5200000, 4100000
    if "RIGHT" in normalized:
        return 6400000, 1600000, 5200000, 4100000
    if "MINI" in normalized:
        return 7800000, 4200000, 2900000, 1700000
    if "SUMMARY" in normalized:
        return 7300000, 1550000, 4000000, 3200000
    return 900000, 1500000, 10300000, 4700000


def replace_placeholder_with_table(
    slide: Any,
    placeholder_name: str,
    csv_path: str | Path,
    *,
    start_row: int = 0,
    max_rows: int | None = None,
    font_size: float = 7.5,
    header_font_size: float = 7.5,
    max_cell_chars: int = 70,
    column_widths: list[float] | None = None,
) -> None:
    shape = _find_placeholder_shape(slide, placeholder_name)
    if shape is None:
        shape = _find_visual_table_placeholder(slide)
    rows = _read_csv_rows(csv_path)
    if not rows:
        raise SlideGenerationError(f"Table CSV is empty: {resolve_path(csv_path)}")
    rows = _slice_table_rows(rows, start_row=start_row, max_rows=max_rows)

    if shape is None:
        left, top, width, height = 571500, 1533525, 11049000, 4750000
    else:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    left, top, width, height = _expand_table_bounds_if_needed(left, top, width, height)
    if shape is not None:
        _remove_shape(shape)
    _add_table_panel(slide, left, top, width, height)
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = table_shape.table
    effective_widths = column_widths if column_widths and len(column_widths) == len(rows[0]) else _column_width_ratios(rows)
    if effective_widths and len(effective_widths) == len(rows[0]):
        for idx, ratio in enumerate(effective_widths):
            table.columns[idx].width = int(width * float(ratio))
    font_size = max(5.6, font_size)
    header_font_size = max(5.8, header_font_size)
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = _compact(value, max_cell_chars)
            _style_dark_table_cell(cell, is_header=row_idx == 0, row_idx=row_idx)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = DECK_FONT_FAMILY
                    run.font.size = _pt(header_font_size if row_idx == 0 else font_size)
                    run.font.bold = row_idx == 0
                    run.font.color.rgb = _rgb("F8FAFC" if row_idx == 0 else "E5E7EB")


def _add_table_panel(slide: Any, left: int, top: int, width: int, height: int) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    pad = 70000
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left - pad, top - pad, width + pad * 2, height + pad * 2)
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb("0B1220")
    panel.fill.transparency = 18
    panel.line.color.rgb = _rgb("334155")
    panel.line.transparency = 15


def _find_visual_table_placeholder(slide: Any) -> Any | None:
    tables = [
        shape
        for shape in slide.shapes
        if shape.shape_type == 19
        and int(getattr(shape, "top", 0)) > 900000
        and int(getattr(shape, "width", 0)) > 1000000
        and int(getattr(shape, "height", 0)) > 1000000
    ]
    return max(tables, key=lambda shape: int(shape.width) * int(shape.height)) if tables else None


def replace_placeholder_with_text_box(slide: Any, placeholder_name: str, text: str) -> None:
    shape = _find_placeholder_shape(slide, placeholder_name)
    if shape is None:
        raise SlideGenerationError(f"Text placeholder not found in slide: {placeholder_name}")
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    _remove_shape(shape)
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = 220000
    box.text_frame.margin_right = 220000
    box.text_frame.margin_top = 160000
    box.text_frame.margin_bottom = 160000
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = _compact(text, 1400)
    for run in paragraph.runs:
        run.font.name = DECK_FONT_FAMILY
        run.font.size = _pt(_fit_font_size(box, paragraph.text, preferred=20, minimum=11))
        run.font.color.rgb = _rgb("F8FAFC")


def format_slide_generation_summary(result: SlideGenerationResult) -> str:
    lines = [
        "Slide generation completed.",
        "",
        f"Template: {result.template}",
        f"Config: {result.config}",
        f"Generated slides: {result.generated_slides}",
        f"Missing optional assets: {len(result.missing_optional_assets)}",
        f"Output: {result.output}",
    ]
    return "\n".join(lines)


def _load_yaml(path: Path) -> dict[str, Any]:
    try:
        data = yaml.safe_load(path.read_text(encoding="utf-8")) or {}
    except yaml.YAMLError as exc:
        raise SlideGenerationError(f"Malformed YAML:\n{path}\n\n{exc}") from exc
    if not isinstance(data, dict):
        raise SlideGenerationError(f"Slide config must contain a YAML mapping:\n{path}")
    return data


def _validate_paths(template: Path, output: Path, overwrite: bool) -> None:
    if not template.exists():
        raise SlideGenerationError(
            f"Template not found:\n{template}\n\nFix:\nPlace the PowerPoint template in `slides/templates/` or pass `--template`."
        )
    if output.exists() and not overwrite:
        backup = output.with_name(f"{output.stem}_{datetime.now().strftime('%Y-%m-%d_%H-%M')}{output.suffix}")
        shutil.copy2(output, backup)


def _template_slides_by_id(presentation: Any) -> dict[str, Any]:
    templates: dict[str, Any] = {}
    for slide in presentation.slides:
        template_id = _extract_template_id(slide)
        if template_id:
            if template_id in templates:
                raise SlideGenerationError(f"Duplicate TEMPLATE_ID in template deck: {template_id}")
            templates[template_id] = slide
    if not templates:
        raise SlideGenerationError("No `TEMPLATE_ID:<id>` markers found in the template deck.")
    return templates


def _extract_template_id(slide: Any) -> str | None:
    for shape in slide.shapes:
        text = _shape_text(shape)
        match = re.search(r"TEMPLATE_ID\s*:\s*([A-Za-z0-9_-]+)", text)
        if match:
            return match.group(1)
    return None


def _validate_requested_slides(slides: list[dict[str, Any]], templates: dict[str, Any]) -> None:
    if not slides:
        raise SlideGenerationError("No slides configured in slide_deck.yml.")
    missing = []
    for spec in slides:
        template_id = _resolved_spec_template_id(spec)
        if template_id not in templates:
            missing.append(template_id)
    if missing:
        available = ", ".join(sorted(templates))
        raise SlideGenerationError(f"Missing TEMPLATE_ID in template: {', '.join(missing)}\nAvailable: {available}")


def _resolved_spec_template_id(spec: dict[str, Any]) -> str:
    base = spec.get("template") or spec.get("template_id")
    theme = spec.get("theme", "neutral")
    try:
        return resolve_template_id(str(base), str(theme))
    except ValueError as exc:
        raise SlideGenerationError(str(exc)) from exc


def validate_template_structure(template_path: str | Path) -> list[str]:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise SlideGenerationError("Missing dependency: python-pptx.") from exc

    path = resolve_path(template_path)
    if not path.exists():
        return [f"Template not found: {path}"]
    messages: list[str] = []
    presentation = Presentation(str(path))
    try:
        template_map = _template_slides_by_id(presentation)
    except SlideGenerationError as exc:
        return [str(exc)]
    for template_id in REQUIRED_TEMPLATE_IDS:
        if template_id not in template_map:
            messages.append(f"Missing TEMPLATE_ID: {template_id}")
    for idx, slide in enumerate(presentation.slides, start=1):
        ids = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            ids.extend(re.findall(r"TEMPLATE_ID\s*:\s*([A-Za-z0-9_-]+)", text))
        if len(ids) != 1:
            messages.append(f"Slide {idx}: expected exactly one TEMPLATE_ID, found {len(ids)}")
    return messages


def _validate_assets(deck_config: dict[str, Any]) -> None:
    errors = _asset_errors(deck_config)
    if errors:
        raise SlideGenerationError(
            "\n\n".join(errors)
            + "\n\nFix:\nRun `python main.py generate-report` first, or update `slides/config/slide_deck.yml`."
        )


def _asset_errors(deck_config: dict[str, Any]) -> list[str]:
    errors: list[str] = []
    for spec in deck_config.get("slides", []):
        for path in (spec.get("images") or {}).values():
            target = resolve_path(path)
            if not target.exists():
                errors.append(f"Missing asset:\n{target}")
        table = spec.get("table")
        if table:
            target = resolve_path(table["source"])
            if not target.exists():
                errors.append(f"Missing table CSV:\n{target}")
        for path in (spec.get("fields_from_file") or {}).values():
            target = resolve_path(path)
            if not target.exists():
                errors.append(f"Missing text file:\n{target}")
    return errors


def _expand_paginated_table_specs(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    expanded: list[dict[str, Any]] = []
    for spec in slides:
        expanded.extend(table_specs_from_paginated_table(spec))
    return expanded


def _collect_fields(spec: dict[str, Any]) -> dict[str, str]:
    fields = {str(key): italian_display_text(value) for key, value in (spec.get("fields") or {}).items()}
    for key, path in (spec.get("fields_from_file") or {}).items():
        fields[str(key)] = italian_display_text(_clean_markdown(resolve_path(path).read_text(encoding="utf-8")))
    return fields


def _replace_text_fields(slide: Any, fields: dict[str, str]) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        shape_name = getattr(shape, "name", "")
        matching_field = _field_for_shape_name(shape_name, fields)
        current_text = _shape_text(shape)
        if matching_field and "TEMPLATE_ID" not in current_text and not _contains_placeholder(current_text, matching_field):
            _set_shape_text_preserving_style(shape, fields[matching_field])
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = _replace_placeholders(run.text, fields)
        _rewrite_markdown_shape_if_needed(shape)


def _replace_named_non_text_fields_with_text(slide: Any, fields: dict[str, str], image_placeholders: set[str]) -> None:
    image_names = {_normalize_placeholder(name) for name in image_placeholders}
    for shape in list(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            continue
        name = getattr(shape, "name", "")
        if not name:
            continue
        normalized = _normalize_placeholder(name)
        if normalized in image_names:
            continue
        for key, value in fields.items():
            if normalized == _normalize_placeholder(key) and str(value).strip():
                _replace_shape_with_text(slide, shape, str(value))
                break


def _replace_shape_with_text(slide: Any, shape: Any, text: str) -> None:
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    _remove_shape(shape)
    box = slide.shapes.add_textbox(left, top, width, height)
    _write_rich_text(box, text, preferred=22, minimum=12, title_like=True)


def _replace_placeholders(text: str, fields: dict[str, str]) -> str:
    result = text
    for key, value in fields.items():
        for variant in _placeholder_variants(key):
            result = result.replace(variant, value)
    return result


def _clear_unresolved_placeholders(slide: Any) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "left", 0) < 0:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = re.sub(r"\[[A-Z0-9_ /:.-]{2,}\]", "", run.text)


def _field_for_shape_name(shape_name: str, fields: dict[str, str]) -> str | None:
    normalized_name = _normalize_placeholder(shape_name)
    for key in fields:
        if normalized_name == _normalize_placeholder(key):
            return key
    return None


def _contains_placeholder(text: str, key: str) -> bool:
    return any(variant in text for variant in _placeholder_variants(key))


def _set_shape_text_preserving_style(shape: Any, text: str) -> None:
    _write_rich_text(shape, text)


def _rewrite_markdown_shape_if_needed(shape: Any) -> None:
    text = _shape_text(shape)
    if "**" in text or re.search(r"(^|\n)\s*[-*]\s+", text) or re.search(r"(^|\n)\s*\d+\.\s+", text):
        _write_rich_text(shape, text)


def _write_rich_text(shape: Any, text: str, *, preferred: int | None = None, minimum: int = 9, title_like: bool = False) -> None:
    from pptx.enum.text import PP_ALIGN

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 130000
    frame.margin_right = 130000
    frame.margin_top = 90000
    frame.margin_bottom = 90000
    lines = _display_lines(text)
    font_size = _fit_font_size(shape, "\n".join(line for _kind, line in lines), preferred=preferred, minimum=minimum) or minimum
    for index, (kind, line) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = ""
        paragraph.space_after = _pt(8 if kind in {"heading", "plain"} else 5)
        paragraph.line_spacing = 1.08
        paragraph.level = 1 if kind == "subbullet" else 0
        paragraph.alignment = PP_ALIGN.LEFT
        if kind in {"bullet", "subbullet"}:
            paragraph.text = "• "
        for segment, bold, italic in _rich_segments(line):
            run = paragraph.add_run()
            run.text = segment
            run.font.name = DECK_FONT_FAMILY
            run.font.size = _pt(font_size + (4 if kind == "heading" else 0))
            run.font.bold = bold or kind == "heading" or title_like
            run.font.italic = italic
            run.font.color.rgb = _rgb(_segment_color(segment, bold=bold, italic=italic, kind=kind))


def _display_lines(text: str) -> list[tuple[str, str]]:
    raw_lines = [line.rstrip() for line in str(text).splitlines()]
    lines: list[tuple[str, str]] = []
    for raw in raw_lines:
        stripped = raw.strip()
        if not stripped:
            continue
        if stripped.startswith("#"):
            lines.append(("heading", stripped.lstrip("#").strip()))
        elif re.match(r"^[-*]\s+", stripped):
            lines.append(("bullet", re.sub(r"^[-*]\s+", "", stripped)))
        elif re.match(r"^\d+\.\s+", stripped):
            lines.append(("bullet", stripped))
        elif raw.startswith("  ") or raw.startswith("\t"):
            lines.append(("subbullet", stripped))
        else:
            lines.append(("plain", stripped))
    return lines or [("plain", "")]


def _rich_segments(text: str) -> list[tuple[str, bool, bool]]:
    pattern = re.compile(r"(\*\*[^*]+\*\*|\*[^*]+\*)")
    segments: list[tuple[str, bool, bool]] = []
    pos = 0
    for match in pattern.finditer(text):
        if match.start() > pos:
            segments.append((text[pos : match.start()], False, False))
        token = match.group(0)
        if token.startswith("**"):
            segments.append((token[2:-2], True, False))
        else:
            segments.append((token[1:-1], False, True))
        pos = match.end()
    if pos < len(text):
        segments.append((text[pos:], False, False))
    return segments or [("", False, False)]


def _segment_color(segment: str, *, bold: bool, italic: bool, kind: str) -> str:
    text = segment.lower()
    if "deliveroo" in text:
        return "00CCBC"
    if "glovo" in text:
        return "FFC244"
    if any(token in text for token in ["ueq", "nps", "user test"]):
        return "A78BFA"
    if any(token in text for token in ["euristica", "priorità", "problemi"]):
        return "38BDF8"
    if kind == "heading" or bold:
        return "F8FAFC"
    if italic:
        return "CBD5E1"
    return "E5E7EB"


def _style_and_fit_slide_text_dark(slide: Any, template_id: str | None = None) -> None:
    is_section = str(template_id or "").startswith("section_divider")
    is_cover = str(template_id or "") == "cover"
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "left", 0) < 0:
            continue
        text = _shape_text(shape).strip()
        if is_cover and _style_cover_title(shape, text):
            continue
        fitted_size = _fit_font_size(shape, text)
        title_size = _single_line_title_size(shape, text)
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = DECK_FONT_FAMILY
                if is_section and text and "SECTION" not in text.upper():
                    shape.text_frame.word_wrap = False
                    run.font.size = _pt(min(_section_title_size(text), title_size or SECTION_TITLE_SIZE_LONG))
                    run.font.bold = True
                    color = _section_brand_color(template_id)
                    if color:
                        run.font.color.rgb = _rgb(color)
                elif title_size and "\n" not in text:
                    shape.text_frame.word_wrap = False
                    current_size = int(run.font.size.pt) if run.font.size else title_size
                    run.font.size = _pt(min(current_size, title_size))
                elif fitted_size:
                    current_size = int(run.font.size.pt) if run.font.size else fitted_size
                    run.font.size = _pt(min(current_size, fitted_size))
                if not (run.font.bold or run.font.italic):
                    run.font.color.rgb = _rgb("F8FAFC")


def _section_brand_color(template_id: str | None) -> str | None:
    template = str(template_id or "").lower()
    if template.endswith("_deliveroo"):
        return DELIVEROO_COLOR
    if template.endswith("_glovo"):
        return GLOVO_COLOR
    return None


def _style_cover_title(shape: Any, text: str) -> bool:
    normalized = " ".join(str(text).split())
    match = re.search(r"deliveroo\s+vs\s+glovo", normalized, flags=re.IGNORECASE)
    if not match:
        return False
    from pptx.enum.text import PP_ALIGN

    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = 0
    frame.margin_right = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    title_size = COVER_TITLE_SIZE if normalized.lower() == "deliveroo vs glovo" else 50
    parts = []
    prefix = normalized[: match.start()]
    suffix = normalized[match.end() :]
    if prefix:
        parts.append((prefix, "F8FAFC"))
    parts.extend([("Deliveroo", DELIVEROO_COLOR), (" vs ", "F8FAFC"), ("Glovo", GLOVO_COLOR)])
    if suffix:
        parts.append((suffix, "F8FAFC"))
    for segment, color in parts:
        run = paragraph.add_run()
        run.text = segment
        run.font.name = DECK_FONT_FAMILY
        run.font.size = _pt(title_size)
        run.font.bold = True
        run.font.color.rgb = _rgb(color)
    return True


def _section_title_size(text: str) -> int:
    length = len(" ".join(str(text).split()))
    if length <= 22:
        return SECTION_TITLE_SIZE_SHORT
    if length <= 36:
        return SECTION_TITLE_SIZE_MEDIUM
    return SECTION_TITLE_SIZE_LONG


def _single_line_title_size(shape: Any, text: str, *, minimum: int = 18) -> int | None:
    compact = " ".join(str(text).split())
    if not compact:
        return None
    width = max(int(getattr(shape, "width", 0)), 1)
    height = max(int(getattr(shape, "height", 0)), 1)
    is_title_like = height <= 950000 and width >= 2500000
    if not is_title_like:
        return None
    max_size = 72 if height >= 650000 else 38
    for size in range(max_size, minimum - 1, -1):
        if len(compact) * size * 6100 <= width * 0.92:
            return size
    return minimum


def _fit_font_size(shape: Any, text: str, *, preferred: int | None = None, minimum: int = 8) -> int | None:
    compact = " ".join(str(text).split())
    if not compact:
        return None
    width = max(int(getattr(shape, "width", 0)), 1)
    height = max(int(getattr(shape, "height", 0)), 1)
    area = width * height
    is_title_like = height <= 750000 and width >= 3000000
    if is_title_like and len(compact) <= 95 and preferred is None:
        return None
    if preferred is not None:
        max_size = preferred
    elif is_title_like:
        max_size = 38
    elif area < 1800000 * 600000:
        max_size = 16
    else:
        max_size = 20
    min_size = minimum if not is_title_like else max(22, minimum)
    for size in range(max_size, min_size - 1, -1):
        chars_per_line = max(8, int(width / max(size * 6200, 1)))
        estimated_lines = 0
        for raw_line in str(text).splitlines() or [compact]:
            line = raw_line.strip() or " "
            estimated_lines += max(1, ceil(len(line) / chars_per_line))
        required_height = estimated_lines * size * 12700 * 1.22
        if required_height <= height * 0.88:
            return size
    return min_size


def _placeholder_variants(key: str) -> list[str]:
    variants = [
        f"[{key}]",
        f"[{key.replace('_', ' ')}]",
        key,
        key.replace("_", " "),
    ]
    aliases = {
        "SUCCESS_RATE": ["[XX%]", "XX%"],
        "SUCCESS_RATE_VALUE": ["[XX%]", "XX%"],
        "AVG_TIME": ["[XX.X s]", "[XX.X S]", "XX.X s", "XX.X S"],
        "AVG_TIME_VALUE": ["[XX.X s]", "[XX.X S]", "XX.X s", "XX.X S"],
    }
    variants.extend(aliases.get(key, []))
    return variants


def _duplicate_slide(presentation: Any, source_slide: Any) -> Any:
    new_slide = presentation.slides.add_slide(_blank_slide_layout(presentation))
    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
            continue
        new_rid = new_slide.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)
        if new_rid != rel.rId:
            _replace_relationship_id(new_slide, rel.rId, new_rid)
    return new_slide


def _blank_slide_layout(presentation: Any) -> Any:
    for layout in presentation.slide_layouts:
        if "blank" in getattr(layout, "name", "").lower():
            return layout
    return presentation.slide_layouts[-1]


def _replace_relationship_id(slide: Any, old_rid: str, new_rid: str) -> None:
    for element in slide.element.iter():
        for attr_name, attr_value in list(element.attrib.items()):
            if attr_value == old_rid:
                element.set(attr_name, new_rid)


def _remove_template_slides(presentation: Any, count: int) -> None:
    slide_ids = list(presentation.slides._sldIdLst)[:count]
    for slide_id in slide_ids:
        presentation.part.drop_rel(slide_id.rId)
        presentation.slides._sldIdLst.remove(slide_id)


def _find_placeholder_shape(slide: Any, placeholder_name: str) -> Any | None:
    normalized = _normalize_placeholder(placeholder_name)
    fallback = None
    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        text = _shape_text(shape)
        if "TEMPLATE_ID" in text:
            continue
        if _normalize_placeholder(name) == normalized or _normalize_placeholder(text) == normalized:
            return shape
        if normalized in _normalize_placeholder(name) or normalized in _normalize_placeholder(text):
            fallback = fallback or shape
    return fallback


def _remove_template_metadata_shapes(slide: Any) -> None:
    for shape in list(slide.shapes):
        if "TEMPLATE_ID" in _shape_text(shape):
            _remove_shape(shape)


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs)


def _normalize_placeholder(value: str) -> str:
    return re.sub(r"[^A-Z0-9]", "", value.upper())


def _remove_shape(shape: Any) -> None:
    shape.element.getparent().remove(shape.element)


def _read_csv_rows(path: str | Path) -> list[list[str]]:
    target = resolve_path(path)
    if not target.exists():
        raise SlideGenerationError(f"Table CSV not found:\n{target}")
    return read_display_table(target)


def _slice_table_rows(rows: list[list[str]], *, start_row: int = 0, max_rows: int | None = None) -> list[list[str]]:
    if not max_rows or max_rows <= 0 or len(rows) <= 1:
        return rows
    header, data = rows[0], rows[1:]
    start = max(0, start_row)
    return [header] + data[start : start + max_rows]


def _expand_table_bounds_if_needed(left: int, top: int, width: int, height: int) -> tuple[int, int, int, int]:
    slide_width = 12192000
    if width >= int(slide_width * 0.7):
        return left, top, width, height
    expanded_left = 571500
    expanded_top = min(top, 1533525)
    expanded_width = 11049000
    expanded_height = 4750000
    return expanded_left, expanded_top, expanded_width, expanded_height


def _column_width_ratios(rows: list[list[str]]) -> list[float]:
    if not rows or not rows[0]:
        return []
    columns = len(rows[0])
    weights = []
    for col_idx in range(columns):
        values = [str(row[col_idx]) for row in rows if col_idx < len(row)]
        max_len = max((len(value) for value in values), default=1)
        avg_len = sum(len(value) for value in values) / max(len(values), 1)
        weights.append(min(max(max_len * 0.55 + avg_len * 0.45, 8), 34))
    total = sum(weights) or 1
    ratios = [weight / total for weight in weights]
    min_ratio = 0.07 if columns <= 8 else 0.055
    ratios = [max(ratio, min_ratio) for ratio in ratios]
    total = sum(ratios) or 1
    return [ratio / total for ratio in ratios]


def _style_dark_table_cell(cell: Any, *, is_header: bool, row_idx: int) -> None:
    cell.fill.solid()
    if is_header:
        cell.fill.fore_color.rgb = _rgb("0B1220")
    elif row_idx % 2:
        cell.fill.fore_color.rgb = _rgb("111827")
    else:
        cell.fill.fore_color.rgb = _rgb("1F2937")
    for side in ("left", "right", "top", "bottom"):
        line = getattr(cell, side, None)
        if line is not None:
            line.color.rgb = _rgb("374151")


def _clean_markdown(text: str) -> str:
    cleaned = re.sub(r"^#{1,6}\s*", "", text, flags=re.MULTILINE)
    cleaned = re.sub(r"[*_`]+", "", cleaned)
    cleaned = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _compact(value: str, limit: int) -> str:
    text = " ".join(str(value).split())
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "..."


def _pt(size: int | float) -> Any:
    from pptx.util import Pt

    return Pt(size)


def _rgb(hex_color: str) -> Any:
    from pptx.dml.color import RGBColor

    value = hex_color.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
